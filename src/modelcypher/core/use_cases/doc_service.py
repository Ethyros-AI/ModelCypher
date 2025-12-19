from __future__ import annotations

import json
import time
import uuid
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

from modelcypher.utils.paths import expand_path


@dataclass(frozen=True)
class DocConvertResult:
    job_id: str
    dataset_name: str
    generator: str
    created_at: str
    duration_seconds: float
    files_processed: int
    sample_count: int
    total_tokens: int
    total_characters: int
    detected_format: str
    output_format: str
    quality_score: int
    validation_status: str
    validation_errors: list[str]
    warnings: list[str]
    source_files: list[str]
    failed_files: list[str]


class DocService:
    def __init__(self) -> None:
        self.supported_extensions = {
            ".md",
            ".txt",
            ".text",
            ".rtf",
            ".swift",
            ".py",
            ".js",
            ".ts",
            ".jsx",
            ".tsx",
            ".java",
            ".kt",
            ".go",
            ".rs",
            ".c",
            ".cpp",
            ".h",
            ".hpp",
            ".cs",
            ".rb",
            ".php",
            ".sh",
            ".bash",
            ".zsh",
            ".fish",
            ".yaml",
            ".yml",
            ".toml",
            ".ini",
            ".cfg",
            ".conf",
            ".json",
            ".jsonl",
            ".csv",
            ".tsv",
            ".xml",
            ".pdf",
            ".html",
            ".htm",
            ".docx",
            ".doc",
            ".pptx",
            ".ppt",
            ".xlsx",
            ".xls",
            ".log",
            ".env",
            ".gitignore",
            ".dockerfile",
            ".makefile",
        }

    def convert(
        self,
        inputs: list[str],
        output_path: str,
        chunk_size: int = 2000,
        chunk_overlap: int = 200,
        text_only: bool = True,
        stream: bool = False,
        update_manifest: bool = False,
    ) -> tuple[DocConvertResult, list[dict]]:
        start = time.time()
        resolved_output = expand_path(output_path)
        resolved_output.parent.mkdir(parents=True, exist_ok=True)

        source_files: list[str] = []
        failed_files: list[str] = []
        warnings: list[str] = []
        events: list[dict] = []

        files = self._collect_files(inputs)
        if stream:
            events.append({"stage": "collecting", "message": f"Found {len(files)} files to process"})

        samples: list[dict] = []
        total_tokens = 0
        total_characters = 0

        for idx, file_path in enumerate(files, start=1):
            if stream:
                events.append({"stage": "loading", "current": idx, "total": len(files), "file": file_path})
            content = self._read_file(file_path)
            if content is None:
                failed_files.append(file_path)
                continue
            source_files.append(file_path)
            total_characters += len(content)
            chunks = self._chunk_text(content, chunk_size, chunk_overlap)
            for chunk in chunks:
                tokens = len(chunk.split())
                total_tokens += tokens
                samples.append({"text": chunk} if text_only else {"messages": [{"role": "user", "content": chunk}]})

        with resolved_output.open("w", encoding="utf-8") as handle:
            for sample in samples:
                handle.write(json.dumps(sample, ensure_ascii=True) + "\n")

        duration = time.time() - start
        result = DocConvertResult(
            job_id=str(uuid.uuid4()),
            dataset_name=Path(output_path).stem,
            generator="ModelCypherCLI",
            created_at=datetime.utcnow().isoformat() + "Z",
            duration_seconds=duration,
            files_processed=len(source_files),
            sample_count=len(samples),
            total_tokens=total_tokens,
            total_characters=total_characters,
            detected_format="text" if text_only else "chat",
            output_format="jsonl",
            quality_score=100,
            validation_status="passed" if samples else "failed",
            validation_errors=[] if samples else ["No samples generated"],
            warnings=warnings,
            source_files=source_files,
            failed_files=failed_files,
        )
        if stream:
            events.append({"stage": "completed", "samples": len(samples), "duration": duration})

        return result, events

    def _collect_files(self, inputs: list[str]) -> list[str]:
        files: list[str] = []
        for entry in inputs:
            path = expand_path(entry)
            if path.is_dir():
                for file_path in path.rglob("*"):
                    if file_path.is_file() and self._is_supported(file_path):
                        files.append(str(file_path))
            elif path.is_file() and self._is_supported(path):
                files.append(str(path))
        return files

    def _is_supported(self, path: Path) -> bool:
        return path.suffix.lower() in self.supported_extensions or path.name.lower() in {
            "makefile",
            "dockerfile",
        }

    def _read_file(self, path: str) -> str | None:
        ext = Path(path).suffix.lower()
        try:
            if ext in {".pdf"}:
                return self._read_pdf(path)
            if ext in {".docx"}:
                return self._read_docx(path)
            if ext in {".pptx"}:
                return self._read_pptx(path)
            if ext in {".xlsx"}:
                return self._read_xlsx(path)
            if ext in {".html", ".htm"}:
                return self._read_html(path)
            return Path(path).read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return None

    @staticmethod
    def _chunk_text(text: str, size: int, overlap: int) -> list[str]:
        if size <= 0:
            return [text]
        chunks = []
        step = max(1, size - overlap)
        for start in range(0, len(text), step):
            chunk = text[start : start + size]
            if chunk:
                chunks.append(chunk)
        return chunks

    @staticmethod
    def _read_pdf(path: str) -> str:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise RuntimeError("pypdf not installed") from exc
        reader = PdfReader(path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    @staticmethod
    def _read_docx(path: str) -> str:
        try:
            import docx
        except ImportError as exc:
            raise RuntimeError("python-docx not installed") from exc
        document = docx.Document(path)
        return "\n".join(paragraph.text for paragraph in document.paragraphs)

    @staticmethod
    def _read_pptx(path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx not installed") from exc
        presentation = Presentation(path)
        texts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    @staticmethod
    def _read_xlsx(path: str) -> str:
        try:
            import openpyxl
        except ImportError as exc:
            raise RuntimeError("openpyxl not installed") from exc
        workbook = openpyxl.load_workbook(path, data_only=True)
        rows = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                rows.append("\t".join(str(cell) if cell is not None else "" for cell in row))
        return "\n".join(rows)

    @staticmethod
    def _read_html(path: str) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise RuntimeError("beautifulsoup4 not installed") from exc
        content = Path(path).read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(content, "html.parser")
        return soup.get_text(separator="\n")
